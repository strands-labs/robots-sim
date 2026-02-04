#!/usr/bin/env python3
"""
Create a multi-page PowerPoint presentation from INTRODUCTION.md for strands-robots-sim
Features dark blue background with white text for professional look
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
import re
import sys
from pathlib import Path


def parse_markdown(md_path):
    """Parse markdown file and extract sections"""
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()

    # Split by ## headings
    sections = []
    lines = content.split('\n')
    current_section = None
    current_content = []

    for line in lines:
        if line.startswith('# '):
            # Main title
            if current_section:
                sections.append({'title': current_section, 'content': '\n'.join(current_content)})
            current_section = line[2:].strip()
            current_content = []
        elif line.startswith('## '):
            # Section heading
            if current_section:
                sections.append({'title': current_section, 'content': '\n'.join(current_content)})
            current_section = line[3:].strip()
            current_content = []
        else:
            current_content.append(line)

    # Add last section
    if current_section:
        sections.append({'title': current_section, 'content': '\n'.join(current_content)})

    return sections


def add_dark_background(slide):
    """Add dark blue background to slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(26, 35, 50)  # Dark blue #1a2332


def add_title_slide(prs, title, subtitle):
    """Add a title slide with dark background and white text"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    add_dark_background(slide)

    # Add title
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(1.5)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title

    p = title_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.name = 'Calibri'
    p.font.color.rgb = RGBColor(255, 255, 255)  # White

    # Add subtitle
    left = Inches(0.5)
    top = Inches(4.0)
    width = Inches(9)
    height = Inches(1)
    subtitle_box = slide.shapes.add_textbox(left, top, width, height)
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle

    p = subtitle_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(22)
    p.font.name = 'Calibri'
    p.font.color.rgb = RGBColor(200, 200, 200)  # Light gray

    return slide


def add_two_column_slide_with_benefits(prs, title, subtitle, left_header, left_content, benefits_content, right_header, right_content, footer=None, right_image_path=None, right_image_caption=None):
    """Add a two-column slide with separate Key Features and Key Benefits boxes

    Args:
        benefits_content: Content for the Key Benefits box (separate from left_content)
        right_image_path: Optional path to image for right column
        right_image_caption: Optional caption text to display below the image
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_dark_background(slide)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.name = 'Calibri'
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.4))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    p = subtitle_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.name = 'Calibri'
    p.font.color.rgb = RGBColor(200, 200, 200)

    # Horizontal separator line under subtitle
    line_h = slide.shapes.add_connector(1, Inches(1.0), Inches(1.7), Inches(9.0), Inches(1.7))
    line_h.line.color.rgb = RGBColor(100, 180, 255)
    line_h.line.width = Pt(1.5)

    # Vertical separator between columns
    line_v = slide.shapes.add_connector(1, Inches(5.0), Inches(2.0), Inches(5.0), Inches(6.6))
    line_v.line.color.rgb = RGBColor(70, 100, 130)
    line_v.line.width = Pt(2)
    line_v.line.dash_style = 2  # Dashed

    # Left column - Key Features background box (top)
    left_bg_features = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.95), Inches(4.3), Inches(2.8)
    )
    left_bg_features.fill.solid()
    left_bg_features.fill.fore_color.rgb = RGBColor(30, 40, 55)
    left_bg_features.line.color.rgb = RGBColor(100, 180, 255)
    left_bg_features.line.width = Pt(1.5)

    # Left column - Key Benefits background box (bottom)
    left_bg_benefits = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(4.85), Inches(4.3), Inches(2.0)
    )
    left_bg_benefits.fill.solid()
    left_bg_benefits.fill.fore_color.rgb = RGBColor(40, 60, 50)
    left_bg_benefits.line.color.rgb = RGBColor(150, 220, 150)
    left_bg_benefits.line.width = Pt(1.5)

    # Right column background box
    right_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(1.95), Inches(4.3), Inches(4.8)
    )
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = RGBColor(30, 40, 55)
    right_bg.line.color.rgb = RGBColor(100, 180, 255)
    right_bg.line.width = Pt(1.5)

    # Left column - Key Features header with underline
    left_header_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.1), Inches(4.0), Inches(0.4))
    lh_frame = left_header_box.text_frame
    lh_frame.text = left_header
    p = lh_frame.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.name = 'Calibri'
    p.font.color.rgb = RGBColor(100, 180, 255)

    # Underline for Key Features header
    line_lh = slide.shapes.add_connector(1, Inches(0.7), Inches(2.55), Inches(4.5), Inches(2.55))
    line_lh.line.color.rgb = RGBColor(100, 180, 255)
    line_lh.line.width = Pt(2)

    # Left column - Key Features content
    left_content_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.7), Inches(4.0), Inches(1.8))
    lc_frame = left_content_box.text_frame
    lc_frame.word_wrap = True
    lc_frame.auto_size = None  # Disable auto-sizing to prevent overflow
    add_formatted_text(lc_frame, left_content)

    # Left column - Key Benefits header
    benefits_header_box = slide.shapes.add_textbox(Inches(0.6), Inches(5.0), Inches(4.0), Inches(0.4))
    bh_frame = benefits_header_box.text_frame
    bh_frame.text = "🎯 Key Benefits"
    p = bh_frame.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.name = 'Calibri'
    p.font.color.rgb = RGBColor(150, 220, 150)

    # Underline for Key Benefits header
    line_bh = slide.shapes.add_connector(1, Inches(0.7), Inches(5.45), Inches(4.5), Inches(5.45))
    line_bh.line.color.rgb = RGBColor(150, 220, 150)
    line_bh.line.width = Pt(2)

    # Left column - Key Benefits content
    benefits_content_box = slide.shapes.add_textbox(Inches(0.7), Inches(5.55), Inches(4.0), Inches(1.2))
    bc_frame = benefits_content_box.text_frame
    bc_frame.word_wrap = True
    bc_frame.auto_size = None  # Disable auto-sizing to prevent overflow
    # Add benefits text as bullet points with smaller font
    lines = benefits_content.strip().split('\n')
    first_para = True
    for line in lines:
        line = line.rstrip()
        if not line:
            continue
        # Keep the bullet point in the text
        text = line.strip()
        if first_para:
            p = bc_frame.paragraphs[0]
            first_para = False
        else:
            p = bc_frame.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(11)
        p.font.name = 'Calibri'
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_after = Pt(3)

    # Right column header with underline
    right_header_box = slide.shapes.add_textbox(Inches(5.3), Inches(2.1), Inches(4.0), Inches(0.4))
    rh_frame = right_header_box.text_frame
    rh_frame.text = right_header
    p = rh_frame.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.name = 'Calibri'
    p.font.color.rgb = RGBColor(100, 180, 255)

    # Underline for right header
    line_rh = slide.shapes.add_connector(1, Inches(5.4), Inches(2.55), Inches(9.2), Inches(2.55))
    line_rh.line.color.rgb = RGBColor(100, 180, 255)
    line_rh.line.width = Pt(2)

    # Right column content - either image or text
    if right_image_path:
        # Add image in right column
        left = Inches(5.4)
        top = Inches(2.7)
        max_width = Inches(3.9)
        max_height = Inches(2.6) if right_image_caption else Inches(3.8)

        try:
            pic = slide.shapes.add_picture(
                right_image_path,
                left, top,
                width=max_width
            )
            if pic.height < max_height:
                pic.top = top + (max_height - pic.height) // 2

            if right_image_caption:
                caption_top = top + max_height + Inches(0.1)
                caption_box = slide.shapes.add_textbox(
                    Inches(5.4), caption_top, Inches(3.9), Inches(0.5)
                )
                caption_frame = caption_box.text_frame
                caption_frame.text = right_image_caption
                caption_frame.word_wrap = True
                p = caption_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(11)
                p.font.italic = True
                p.font.name = 'Calibri'
                p.font.color.rgb = RGBColor(180, 200, 220)
        except Exception as e:
            print(f"⚠️ Could not load image {right_image_path}: {e}")
            right_content_box = slide.shapes.add_textbox(Inches(5.4), Inches(2.7), Inches(4.0), Inches(3.9))
            rc_frame = right_content_box.text_frame
            rc_frame.word_wrap = True
            add_formatted_text(rc_frame, right_content)
    else:
        right_content_box = slide.shapes.add_textbox(Inches(5.4), Inches(2.7), Inches(4.0), Inches(3.9))
        rc_frame = right_content_box.text_frame
        rc_frame.word_wrap = True
        add_formatted_text(rc_frame, right_content)

    # Footer with background
    if footer:
        footer_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(6.9), Inches(9), Inches(0.5)
        )
        footer_bg.fill.solid()
        footer_bg.fill.fore_color.rgb = RGBColor(30, 40, 55)
        footer_bg.line.color.rgb = RGBColor(150, 220, 150)
        footer_bg.line.width = Pt(1.5)

        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.95), Inches(9), Inches(0.4))
        footer_frame = footer_box.text_frame
        footer_frame.text = footer
        p = footer_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.name = 'Consolas'
        p.font.color.rgb = RGBColor(150, 220, 150)

    return slide


def add_two_column_slide(prs, title, subtitle, left_header, left_content, right_header, right_content, footer=None, right_image_path=None, right_image_caption=None):
    """Add a two-column slide with headers and visual separators

    Args:
        right_image_path: Optional path to image for right column (replaces right_content text)
        right_image_caption: Optional caption text to display below the image
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_dark_background(slide)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.name = 'Calibri'
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.4))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    p = subtitle_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.name = 'Calibri'
    p.font.color.rgb = RGBColor(200, 200, 200)

    # Horizontal separator line under subtitle
    line_h = slide.shapes.add_connector(1, Inches(1.0), Inches(1.7), Inches(9.0), Inches(1.7))
    line_h.line.color.rgb = RGBColor(100, 180, 255)
    line_h.line.width = Pt(1.5)

    # Vertical separator between columns
    line_v = slide.shapes.add_connector(1, Inches(5.0), Inches(2.0), Inches(5.0), Inches(6.6))
    line_v.line.color.rgb = RGBColor(70, 100, 130)
    line_v.line.width = Pt(2)
    line_v.line.dash_style = 2  # Dashed

    # Left column background box
    left_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.95), Inches(4.3), Inches(4.8)
    )
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = RGBColor(30, 40, 55)
    left_bg.line.color.rgb = RGBColor(100, 180, 255)
    left_bg.line.width = Pt(1.5)

    # Right column background box
    right_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(1.95), Inches(4.3), Inches(4.8)
    )
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = RGBColor(30, 40, 55)
    right_bg.line.color.rgb = RGBColor(100, 180, 255)
    right_bg.line.width = Pt(1.5)

    # Left column header with underline
    left_header_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.1), Inches(4.0), Inches(0.4))
    lh_frame = left_header_box.text_frame
    lh_frame.text = left_header
    p = lh_frame.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.name = 'Calibri'
    p.font.color.rgb = RGBColor(100, 180, 255)

    # Underline for left header
    line_lh = slide.shapes.add_connector(1, Inches(0.7), Inches(2.55), Inches(4.5), Inches(2.55))
    line_lh.line.color.rgb = RGBColor(100, 180, 255)
    line_lh.line.width = Pt(2)

    # Left column content
    left_content_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.7), Inches(4.0), Inches(3.9))
    lc_frame = left_content_box.text_frame
    lc_frame.word_wrap = True
    add_formatted_text(lc_frame, left_content)

    # Right column header with underline
    right_header_box = slide.shapes.add_textbox(Inches(5.3), Inches(2.1), Inches(4.0), Inches(0.4))
    rh_frame = right_header_box.text_frame
    rh_frame.text = right_header
    p = rh_frame.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.name = 'Calibri'
    p.font.color.rgb = RGBColor(100, 180, 255)

    # Underline for right header
    line_rh = slide.shapes.add_connector(1, Inches(5.4), Inches(2.55), Inches(9.2), Inches(2.55))
    line_rh.line.color.rgb = RGBColor(100, 180, 255)
    line_rh.line.width = Pt(2)

    # Right column content - either image or text
    if right_image_path:
        # Add image in right column
        left = Inches(5.4)
        top = Inches(2.7)
        # Calculate dimensions to fit within the box
        max_width = Inches(3.9)
        # Reserve space for caption if provided
        max_height = Inches(2.6) if right_image_caption else Inches(3.8)

        try:
            pic = slide.shapes.add_picture(
                right_image_path,
                left, top,
                width=max_width
            )
            # Center the image vertically if needed
            if pic.height < max_height:
                pic.top = top + (max_height - pic.height) // 2

            # Add caption below image if provided
            if right_image_caption:
                caption_top = top + max_height + Inches(0.1)
                caption_box = slide.shapes.add_textbox(
                    Inches(5.4), caption_top, Inches(3.9), Inches(0.5)
                )
                caption_frame = caption_box.text_frame
                caption_frame.text = right_image_caption
                caption_frame.word_wrap = True
                p = caption_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(11)
                p.font.italic = True
                p.font.name = 'Calibri'
                p.font.color.rgb = RGBColor(180, 200, 220)
        except Exception as e:
            # Fallback to text if image fails to load
            print(f"⚠️ Could not load image {right_image_path}: {e}")
            right_content_box = slide.shapes.add_textbox(Inches(5.4), Inches(2.7), Inches(4.0), Inches(3.9))
            rc_frame = right_content_box.text_frame
            rc_frame.word_wrap = True
            add_formatted_text(rc_frame, right_content)
    else:
        # Add text in right column
        right_content_box = slide.shapes.add_textbox(Inches(5.4), Inches(2.7), Inches(4.0), Inches(3.9))
        rc_frame = right_content_box.text_frame
        rc_frame.word_wrap = True
        add_formatted_text(rc_frame, right_content)

    # Footer with background
    if footer:
        footer_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(6.9), Inches(9), Inches(0.5)
        )
        footer_bg.fill.solid()
        footer_bg.fill.fore_color.rgb = RGBColor(30, 40, 55)
        footer_bg.line.color.rgb = RGBColor(150, 220, 150)
        footer_bg.line.width = Pt(1.5)

        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.95), Inches(9), Inches(0.4))
        footer_frame = footer_box.text_frame
        footer_frame.text = footer
        p = footer_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.name = 'Consolas'
        p.font.color.rgb = RGBColor(150, 220, 150)

    return slide


def add_table_slide(prs, title, headers, rows):
    """Add a slide with a comparison table"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_dark_background(slide)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = 'Calibri'
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Add table
    rows_count = len(rows) + 1  # +1 for header
    cols_count = len(headers)
    left = Inches(1.0)
    top = Inches(1.5)
    width = Inches(8.0)
    height = Inches(5.0)

    table = slide.shapes.add_table(rows_count, cols_count, left, top, width, height).table

    # Set column widths
    for col_idx in range(cols_count):
        table.columns[col_idx].width = Inches(8.0 / cols_count)

    # Header row
    for col_idx, header in enumerate(headers):
        cell = table.rows[0].cells[col_idx]
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(50, 100, 150)  # Medium blue
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.rows[row_idx + 1].cells[col_idx]
            cell.text = cell_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(35, 45, 60)  # Slightly lighter than background
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = RGBColor(255, 255, 255)

    return slide


def add_bar_chart_slide(prs, title, categories, series_data):
    """Add a slide with a bar chart"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_dark_background(slide)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = 'Calibri'
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Add chart
    chart_data = CategoryChartData()
    chart_data.categories = categories

    for series_name, values in series_data.items():
        chart_data.add_series(series_name, values)

    x, y, cx, cy = Inches(1.5), Inches(1.5), Inches(7.0), Inches(5.0)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    # Style the chart
    plot = chart.plots[0]
    plot.has_data_labels = True

    return slide


def add_architecture_diagram_slide(prs, title):
    """Add a slide with enhanced architecture flow diagram integrating System 1 and System 2"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_dark_background(slide)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.name = 'Calibri'
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # System 2 Container (large background box)
    sys2_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.6), Inches(1.0), Inches(8.8), Inches(2.3)
    )
    sys2_bg.fill.solid()
    sys2_bg.fill.fore_color.rgb = RGBColor(40, 60, 90)
    sys2_bg.line.color.rgb = RGBColor(100, 150, 255)
    sys2_bg.line.width = Pt(3)

    # System 2 Label
    sys2_label = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(2.5), Inches(0.4))
    sys2_label_text = sys2_label.text_frame
    sys2_label_text.text = "🤔 SYSTEM 2: Deliberate Planning"
    p = sys2_label_text.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(150, 200, 255)

    # LLM Agent box
    agent_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2.0), Inches(1.7), Inches(6.0), Inches(1.2)
    )
    agent_box.fill.solid()
    agent_box.fill.fore_color.rgb = RGBColor(80, 120, 180)
    agent_box.line.color.rgb = RGBColor(150, 200, 255)
    agent_box.line.width = Pt(2.5)
    agent_text = agent_box.text_frame
    agent_text.text = "🤖 Strands Agent (Claude LLM)\nTask Reasoning • Language Understanding • Planning"
    p = agent_text.paragraphs[0]
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Arrow from System 2 to System 1
    arrow_s2_s1 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(4.7), Inches(3.3), Inches(0.6), Inches(0.5))
    arrow_s2_s1.fill.solid()
    arrow_s2_s1.fill.fore_color.rgb = RGBColor(150, 200, 255)
    arrow_s2_s1.line.fill.background()

    # Language instruction label
    instr_label = slide.shapes.add_textbox(Inches(3.0), Inches(3.5), Inches(4.0), Inches(0.3))
    instr_text = instr_label.text_frame
    instr_text.text = "Language Instructions"
    p = instr_text.paragraphs[0]
    p.font.size = Pt(11)
    p.font.italic = True
    p.font.color.rgb = RGBColor(200, 220, 255)
    p.alignment = PP_ALIGN.CENTER

    # System 1 Container (large background box)
    sys1_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.6), Inches(3.9), Inches(8.8), Inches(2.9)
    )
    sys1_bg.fill.solid()
    sys1_bg.fill.fore_color.rgb = RGBColor(50, 70, 50)
    sys1_bg.line.color.rgb = RGBColor(150, 220, 150)
    sys1_bg.line.width = Pt(3)

    # System 1 Label
    sys1_label = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(2.5), Inches(0.4))
    sys1_label_text = sys1_label.text_frame
    sys1_label_text.text = "⚡ SYSTEM 1: Fast Execution"
    p = sys1_label_text.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(150, 220, 150)

    # GR00T VLA box
    groot_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.5), Inches(4.6), Inches(3.5), Inches(0.9)
    )
    groot_box.fill.solid()
    groot_box.fill.fore_color.rgb = RGBColor(70, 130, 70)
    groot_box.line.color.rgb = RGBColor(150, 220, 150)
    groot_box.line.width = Pt(2.5)
    groot_text = groot_box.text_frame
    groot_text.text = "🧠 GR00T VLA Policy\nVision-Language-Action (3B)"
    p = groot_text.paragraphs[0]
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # SimEnv box
    simenv_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.0), Inches(4.6), Inches(3.5), Inches(0.9)
    )
    simenv_box.fill.solid()
    simenv_box.fill.fore_color.rgb = RGBColor(70, 100, 130)
    simenv_box.line.color.rgb = RGBColor(100, 180, 255)
    simenv_box.line.width = Pt(2.5)
    simenv_text = simenv_box.text_frame
    simenv_text.text = "⚙️ SimEnv/SteppedSimEnv\nExecution Engine"
    p = simenv_text.paragraphs[0]
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Libero Simulation box
    libero_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.0), Inches(5.7), Inches(3.5), Inches(0.9)
    )
    libero_box.fill.solid()
    libero_box.fill.fore_color.rgb = RGBColor(60, 90, 120)
    libero_box.line.color.rgb = RGBColor(100, 180, 255)
    libero_box.line.width = Pt(2.5)
    libero_text = libero_box.text_frame
    libero_text.text = "🎮 Libero Simulation\n90+ Benchmark Tasks"
    p = libero_text.paragraphs[0]
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # SimEnv <-> GR00T bidirectional arrows
    # Right arrow (instruction)
    arrow_r = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.5), Inches(4.95), Inches(1.0), Inches(0.3))
    arrow_r.fill.solid()
    arrow_r.fill.fore_color.rgb = RGBColor(255, 200, 100)
    arrow_r.line.fill.background()

    # Instruction label on arrow
    instr_label_r = slide.shapes.add_textbox(Inches(4.6), Inches(4.98), Inches(0.8), Inches(0.24))
    instr_text_r = instr_label_r.text_frame
    instr_text_r.text = "Instruction"
    p = instr_text_r.paragraphs[0]
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Left arrow (action)
    arrow_l = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(4.5), Inches(5.3), Inches(1.0), Inches(0.3))
    arrow_l.fill.solid()
    arrow_l.fill.fore_color.rgb = RGBColor(150, 220, 150)
    arrow_l.line.fill.background()

    # Action label on arrow
    action_label_l = slide.shapes.add_textbox(Inches(4.7), Inches(5.33), Inches(0.6), Inches(0.24))
    action_text_l = action_label_l.text_frame
    action_text_l.text = "Action"
    p = action_text_l.paragraphs[0]
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # ZMQ label
    zmq_label = slide.shapes.add_textbox(Inches(4.7), Inches(4.6), Inches(0.6), Inches(0.3))
    zmq_text = zmq_label.text_frame
    zmq_text.text = "ZMQ"
    p = zmq_text.paragraphs[0]
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

    # SimEnv <-> Libero arrows
    # Down arrow
    arrow_d = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(2.4), Inches(5.5), Inches(0.5), Inches(0.2))
    arrow_d.fill.solid()
    arrow_d.fill.fore_color.rgb = RGBColor(100, 180, 255)
    arrow_d.line.fill.background()

    # Up arrow (feedback)
    arrow_u = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(3.1), Inches(5.5), Inches(0.5), Inches(0.2))
    arrow_u.fill.solid()
    arrow_u.fill.fore_color.rgb = RGBColor(255, 200, 100)
    arrow_u.line.fill.background()

    return slide


def add_formatted_text(text_frame, content):
    """Add formatted text with bullet points to a text frame"""
    lines = content.strip().split('\n')
    first_para = True

    for line in lines:
        line = line.rstrip()
        if not line:
            continue

        # Remove markdown formatting
        clean_text = re.sub(r'\*\*(.*?)\*\*', r'\1', line)
        clean_text = re.sub(r'\*(.*?)\*', r'\1', clean_text)
        clean_text = re.sub(r'`(.*?)`', r'\1', clean_text)

        # Determine indentation level
        if line.startswith('  '):
            level = 1
            text = clean_text.lstrip('- ').strip()
        else:
            level = 0
            text = clean_text.lstrip('- ').strip()

        if first_para:
            p = text_frame.paragraphs[0]
            first_para = False
        else:
            p = text_frame.add_paragraph()

        p.text = text
        p.level = level

        # Check if this is a command line (starts with bash, python, conda, etc.)
        is_command = any(text.startswith(cmd) for cmd in ['bash ', 'python ', 'conda ', 'pip ', 'git ', 'npm ', 'docker '])

        if is_command:
            # Style as command (green, monospace)
            p.font.size = Pt(12)
            p.font.name = 'Consolas'
            p.font.color.rgb = RGBColor(150, 220, 150)  # Light green
        else:
            # Normal text styling
            p.font.size = Pt(13) if level == 0 else Pt(11)
            p.font.name = 'Calibri'
            p.font.color.rgb = RGBColor(255, 255, 255) if level == 0 else RGBColor(200, 200, 200)

        p.space_after = Pt(4)


def add_content_slide(prs, title, content):
    """Add a content slide with dark background and white text"""
    # Use blank layout for custom styling
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_dark_background(slide)

    # Add title with styling
    left = Inches(0.5)
    top = Inches(0.4)
    width = Inches(9)
    height = Inches(0.7)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title

    p = title_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = 'Calibri'
    p.font.color.rgb = RGBColor(255, 255, 255)  # White

    # Add content area
    left = Inches(0.7)
    top = Inches(1.3)
    width = Inches(8.6)
    height = Inches(5.8)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    # Parse content for bullet points and code blocks
    lines = content.strip().split('\n')
    in_code_block = False
    code_lines = []

    for line in lines:
        line = line.rstrip()

        # Handle code blocks
        if line.startswith('```'):
            if in_code_block:
                # End of code block
                if code_lines:
                    code_text = '\n'.join(code_lines)
                    p = text_frame.add_paragraph()
                    p.text = code_text
                    p.font.name = 'Consolas'
                    p.font.size = Pt(11)
                    p.font.color.rgb = RGBColor(150, 220, 150)  # Light green for code on dark bg
                    p.level = 0
                    p.space_after = Pt(12)
                    code_lines = []
                in_code_block = False
            else:
                in_code_block = True
            continue

        if in_code_block:
            code_lines.append(line)
            continue

        # Skip empty lines at the start
        if not text_frame.paragraphs[0].text and not line.strip():
            continue

        # Handle bullet points and regular text
        if line.startswith('- ') or line.startswith('* '):
            # Bullet point
            bullet_text = line[2:].strip()
            # Remove markdown formatting
            bullet_text = re.sub(r'\*\*(.*?)\*\*', r'\1', bullet_text)
            bullet_text = re.sub(r'\*(.*?)\*', r'\1', bullet_text)
            bullet_text = re.sub(r'`(.*?)`', r'\1', bullet_text)

            if not text_frame.paragraphs[0].text:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = bullet_text
            p.level = 0
            p.font.size = Pt(16)
            p.font.name = 'Calibri'
            p.font.color.rgb = RGBColor(255, 255, 255)  # White
            p.space_after = Pt(8)
        elif line.startswith('  - ') or line.startswith('  * '):
            # Sub-bullet
            bullet_text = line[4:].strip()
            bullet_text = re.sub(r'\*\*(.*?)\*\*', r'\1', bullet_text)
            bullet_text = re.sub(r'\*(.*?)\*', r'\1', bullet_text)
            bullet_text = re.sub(r'`(.*?)`', r'\1', bullet_text)

            p = text_frame.add_paragraph()
            p.text = bullet_text
            p.level = 1
            p.font.size = Pt(14)
            p.font.name = 'Calibri'
            p.font.color.rgb = RGBColor(200, 200, 200)  # Light gray
            p.space_after = Pt(6)
        elif line.startswith('|') or line.startswith('---'):
            # Skip markdown table separators for now
            continue
        elif line.strip() and not line.startswith('#'):
            # Regular text
            clean_text = re.sub(r'\*\*(.*?)\*\*', r'\1', line)
            clean_text = re.sub(r'\*(.*?)\*', r'\1', clean_text)
            clean_text = re.sub(r'`(.*?)`', r'\1', clean_text)
            clean_text = re.sub(r'\[(.*?)\]\(.*?\)', r'\1', clean_text)

            if clean_text.strip():
                if not text_frame.paragraphs[0].text:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = clean_text
                p.level = 0
                p.font.size = Pt(16)
                p.font.name = 'Calibri'
                p.font.color.rgb = RGBColor(255, 255, 255)  # White
                p.space_after = Pt(8)

    return slide


def create_presentation_from_markdown(md_path, output_path):
    """Create presentation from markdown file"""
    # Create a new blank presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Parse markdown
    sections = parse_markdown(md_path)

    if not sections:
        print("No sections found in markdown file")
        return

    # Create section dictionary for easy access
    section_dict = {s['title']: s['content'] for s in sections}

    # Get base directory from markdown path
    base_dir = Path(md_path).parent.parent

    # Slide 1: Overview slide with two columns and separate benefits box
    video_path = base_dir / "images" / "libero_example.mp4"
    add_two_column_slide_with_benefits(
        prs,
        title="Strands Robots Simulation",
        subtitle="Robot Control for Strands Agents in Simulated Environments",
        left_header="✨ Key Features",
        left_content="""• Two Execution Modes
  - SimEnv: Full episode
  - SteppedSimEnv: Iterative control

• GR00T VLA Policy (N1.5-3B via ZMQ)

• Libero Simulation (90+ tasks)

• Video Recording & Monitoring""",
        benefits_content="""• Rapid prototyping without physical hardware
• Safe, simulated testing environment
• Iterate on agent strategies & VLA policies
• Validate approaches before real-world deployment""",
        right_header="🎬 Example in Action",
        right_content="",  # Will be replaced by image
        footer="🚀 Quick Start: python examples/libero_example.py | libero_stepped_example.py",
        right_image_path=str(video_path),
        right_image_caption="Robot putting frying pan on cabinet shelf in Libero-90 task suite"
    )

    # Slide 2: How it Works - Flow diagram with image
    flow_diagram_path = base_dir / "images" / "how_it_works_flow.png"
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_dark_background(slide)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "💡 How it Works"
    p = title_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.name = 'Calibri'
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(9), Inches(0.35))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "A streamlined pipeline enabling natural language control of simulated robots"
    p = subtitle_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.name = 'Calibri'
    p.font.color.rgb = RGBColor(200, 200, 200)

    # Add flow diagram image
    try:
        # Flow diagram is now vertical, constrain by height and center horizontally
        pic = slide.shapes.add_picture(
            str(flow_diagram_path),
            Inches(0.5), Inches(1.5),
            height=Inches(5)
        )
        # Center horizontally
        available_width = Inches(9)
        if pic.width < available_width:
            pic.left = Inches(0.5) + (available_width - pic.width) // 2
    except Exception as e:
        print(f"⚠️ Could not load flow diagram image: {e}")
        # Fallback to text
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        text_frame.text = "Flow diagram could not be loaded. Please regenerate the image."
        p = text_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)

    # Footer note
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Action chunks (8 steps) transmitted to Libero environment sequentially • Visual feedback flows back to agent"
    p = footer_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.name = 'Calibri'
    p.font.color.rgb = RGBColor(180, 200, 220)

    # Slide 3: System 1 vs System 2 Thinking
    add_two_column_slide(
        prs,
        title="🧠 System 1 vs System 2 Thinking",
        subtitle="Dual-system architecture inspired by cognitive science",
        left_header="🤔 System 2: Deliberate Planning",
        left_content="""Strands Agent (Claude LLM)


• Task reasoning & decomposition
• Natural language understanding
• Strategy adaptation
• Error recovery


Provides:
✓ Complex reasoning
✓ Contextual understanding
✓ Adaptive planning""",
        right_header="⚡ System 1: Fast Execution",
        right_content="""GR00T VLA Policy


• Vision + Language → Actions
• Low-level control
• Real-time feedback


Provides:
✓ 40-160ms inference
✓ Visuomotor skills
✓ Reactive execution
✓ Task generalization"""
    )

    # Slide 4: Architecture Overview with image
    arch_diagram_path = base_dir / "images" / "architecture_overview.png"
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_dark_background(slide)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "📐 Architecture Overview"
    p = title_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.name = 'Calibri'
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(9), Inches(0.35))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Five interconnected layers enabling natural language robot control"
    p = subtitle_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.name = 'Calibri'
    p.font.color.rgb = RGBColor(200, 200, 200)

    # Add architecture diagram image
    try:
        # Architecture diagram is vertical, so constrain by height
        pic = slide.shapes.add_picture(
            str(arch_diagram_path),
            Inches(0.5), Inches(1.5),
            height=Inches(5)
        )
        # Center horizontally if needed
        available_width = Inches(9)
        if pic.width < available_width:
            pic.left = Inches(0.5) + (available_width - pic.width) // 2
    except Exception as e:
        print(f"⚠️ Could not load architecture diagram image: {e}")
        # Fallback to text
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        text_frame.text = "Architecture diagram could not be loaded. Please regenerate the image."
        p = text_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)

    # Footer note
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Modular design enables swapping policy implementations or simulation environments without restructuring core logic"
    p = footer_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.name = 'Calibri'
    p.font.color.rgb = RGBColor(180, 200, 220)

    # Slide 5: Mode Comparison Table
    add_table_slide(
        prs,
        "⚖️ SimEnv vs SteppedSimEnv Comparison",
        headers=["Feature", "SimEnv", "SteppedSimEnv"],
        rows=[
            ["Control Flow", "One-shot execution", "Step-by-step iteration"],
            ["Agent Feedback", "Final reward only", "Visual + state per batch"],
            ["Use Case", "Known tasks", "Complex tasks"],
            ["Execution Speed", "⚡ Faster", "🐢 Slower (more control)"],
            ["Error Recovery", "None", "✓ Adaptive retry"],
            ["Visual Grounding", "None", "✓ Camera feedback"],
            ["Debugging", "Limited", "✓ Full trace"],
        ]
    )

    # Slide 6: Quick Start
    add_two_column_slide(
        prs,
        title="🚀 Quick Start",
        subtitle="Setup and run your first robot simulation",
        left_header="📋 Prerequisites & Execution",
        left_content="""1. Setup Isaac-GR00T:

bash scripts/setup-gr00t-gpu.sh


2. Install environment:

conda env create -f environment.yml


3. Run Stepped Example:

python examples/libero_stepped_example.py


4. Run Standard Example:

python examples/libero_example.py""",
        right_header="📁 Results in ./rollouts/",
        right_content="""Videos saved automatically:


rollouts/
├── 2026_01_13/
│   ├── 2026_01_13_14_30_45
│       --episode=0
│       --success=True.mp4
│   └── 2026_01_13_14_32_10
│       --episode=1
│       --success=False.mp4
└── 2026_01_12/
    └── ...


• Dual camera side-by-side view
• 30 FPS, full episode"""
    )

    # Slide 7: Supported Environments
    add_two_column_slide(
        prs,
        title="🌍 Supported Environments",
        subtitle="Simulation environments and VLA policies",
        left_header="✅ Current Support",
        left_content="""Libero Task Suites:
• libero_spatial (10 tasks)
• libero_object (10 tasks)
• libero_goal (10 tasks)
• libero_10 (10 tasks)
• libero_90 (90 tasks)


GR00T VLA Policy:
• Isaac-GR00T N1.5-3B
• Checkpoints: spatial, 90-task
• Communication via ZMQ""",
        right_header="🚀 Coming Soon",
        right_content="""Extended Environments:


• IsaacLab


Additional VLA Policies:


• ACT (Action Chunking)


• SmolVLA


• Custom VLA models"""
    )

    # Save presentation
    prs.save(output_path)
    print(f"✅ Multi-page PowerPoint presentation created: {output_path}")
    print(f"   Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    # Paths
    base_dir = Path(__file__).parent.parent
    md_path = base_dir / "docs" / "INTRODUCTION.md"
    output_path = base_dir / "docs" / "Strands_Robots_Simulation.pptx"

    # Check if markdown file exists
    if not md_path.exists():
        print(f"❌ Markdown file not found: {md_path}")
        sys.exit(1)

    # Create presentation
    create_presentation_from_markdown(md_path, output_path)
